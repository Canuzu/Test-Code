# -*- coding: utf-8 -*-
"""Build the RoyalCards pitch deck (16:9) — on-brand with the live site.

Palette: Cream #FAF6EF · Royal Violet #5B2C91 · Ink #2C1A4D · Gold #C9A961.
Fonts:   Cormorant Garamond (display), Great Vibes (script logo), Inter (body).
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

HERE = os.path.dirname(__file__)
ASSETS = os.path.join(HERE, "assets")
OUT = os.path.join(os.path.dirname(HERE), "RoyalCards_Pitch_Deck.pptx")

# ----------------------------------------------------------------------------
# Palette
INK = "2C1A4D"
INK_SOFT = "4A3472"
VIOLET = "5B2C91"
VIOLET_DEEP = "3A1B5C"
VIOLET_LIGHT = "8B5CC4"
GOLD = "C9A961"
GOLD_DARK = "9A7B3C"
GOLD_LIGHT = "D9C089"
CREAM = "FAF6EF"
CREAM_WARM = "F5EFE0"
CREAM_DEEP = "E8DCC4"
WHITE = "FFFFFF"
PANEL = "FCFAF5"
BORDER = "E7DCC6"
BORDER_V = "E4DAF0"
MUTED = "8A7BA0"
# on violet slides
ON_V_TEXT = "F0EAFF"
ON_V_MUTED = "C4B5E8"
ON_V_PANEL = "2A1A47"
ON_V_BORDER = "5A3E86"

# Fonts
DISPLAY = "Cormorant Garamond"
SCRIPT = "Great Vibes"
BODY = "Inter"
SYM = "Segoe UI Symbol"

# Geometry (inches)
SW, SH = 13.333, 7.5
M = 0.92
CW = SW - 2 * M
TOTAL = 13

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def C(hexstr):
    return RGBColor.from_string(hexstr)


def slide_bg(img):
    s = prs.slides.add_slide(BLANK)
    s.shapes.add_picture(os.path.join(ASSETS, img), 0, 0,
                         width=prs.slide_width, height=prs.slide_height)
    return s


def _set_spacing(run, val):
    """Letter spacing in 1/100 pt."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(val)))


def text(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True, fit=True):
    """paras: list of dicts {align,line,space_before,space_after,runs:[...]}.
    run dict: {t, font, s, c, b, i, sp}."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    # shrink-to-fit safety net: keeps text inside its box across font metrics
    if fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        if "line" in para:
            p.line_spacing = para["line"]
        if "space_before" in para:
            p.space_before = Pt(para["space_before"])
        p.space_after = Pt(para.get("space_after", 0))
        for r in para["runs"]:
            run = p.add_run()
            run.text = r["t"]
            f = run.font
            f.name = r.get("font", BODY)
            f.size = Pt(r["s"])
            f.bold = r.get("b", False)
            f.italic = r.get("i", False)
            f.color.rgb = C(r.get("c", INK))
            if r.get("sp"):
                _set_spacing(run, r["sp"])
    return tb


def no_shadow(shape):
    shape.shadow.inherit = False


def soft_shadow(shape, blur=0.09, dist=0.055, direction=5400000,
                color=INK, opacity=22):
    """Add a subtle outer shadow via XML. opacity in percent (alpha)."""
    spPr = shape._element.spPr
    for old in spPr.findall(qn("a:effectLst")):
        spPr.remove(old)
    eff = spPr.makeelement(qn("a:effectLst"), {})
    sh = eff.makeelement(qn("a:outerShdw"), {
        "blurRad": str(int(blur * 914400)),
        "dist": str(int(dist * 914400)),
        "dir": str(direction),
        "rotWithShape": "0",
    })
    clr = sh.makeelement(qn("a:srgbClr"), {"val": color})
    al = clr.makeelement(qn("a:alpha"), {"val": str(int(opacity * 1000))})
    clr.append(al)
    sh.append(clr)
    eff.append(sh)
    spPr.append(eff)


def rrect(slide, x, y, w, h, fill=PANEL, line=BORDER, line_w=1.0,
          radius=0.055, shadow=True, shadow_kw=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = radius
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = C(fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = C(line)
        shp.line.width = Pt(line_w)
    no_shadow(shp)
    if shadow:
        soft_shadow(shp, **(shadow_kw or {}))
    return shp


def rect(slide, x, y, w, h, fill, line=None, line_w=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = C(fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = C(line)
        shp.line.width = Pt(line_w)
    no_shadow(shp)
    return shp


def oval(slide, x, y, w, h, fill, line=None, line_w=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = C(fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = C(line)
        shp.line.width = Pt(line_w)
    no_shadow(shp)
    return shp


def hline(slide, x, y, w, color=GOLD, weight=1.4):
    """Thin horizontal rule using a connector."""
    ln = slide.shapes.add_connector(2, Inches(x), Inches(y),
                                    Inches(x + w), Inches(y))
    ln.line.color.rgb = C(color)
    ln.line.width = Pt(weight)
    no_shadow(ln)
    return ln


def eyebrow(slide, x, y, txt, color=VIOLET, w=6.0, align=PP_ALIGN.LEFT):
    text(slide, x, y, w, 0.32, [{
        "align": align,
        "runs": [{"t": "—  " + txt, "font": BODY, "s": 11.5, "c": color,
                  "b": True, "sp": 290}],
    }])


def footer(slide, n, dark=False):
    wm_c = GOLD if dark else VIOLET
    pg_c = ON_V_MUTED if dark else MUTED
    rule_c = "5A3E86" if dark else CREAM_DEEP
    hline(slide, M, 7.04, CW, color=rule_c, weight=0.75)
    text(slide, M, 7.08, 4.0, 0.3, [{
        "runs": [{"t": "RoyalCards", "font": SCRIPT, "s": 17, "c": wm_c}],
    }])
    text(slide, SW - M - 4.0, 7.12, 4.0, 0.3, [{
        "align": PP_ALIGN.RIGHT,
        "runs": [
            {"t": "PREMIUM TCG  ·  DÜSSELDORF      ", "font": BODY, "s": 8,
             "c": pg_c, "sp": 120},
            {"t": "%02d / %02d" % (n, TOTAL), "font": BODY, "s": 8.5,
             "c": (GOLD if dark else GOLD_DARK), "b": True, "sp": 120},
        ],
    }])


def crown(slide, x, y, size=30, color=GOLD, align=PP_ALIGN.LEFT, w=1.2):
    text(slide, x, y, w, size / 50.0 + 0.5, [{
        "align": align,
        "runs": [{"t": "♛", "font": SYM, "s": size, "c": color}],
    }], fit=False)


def header(slide, eb, title_runs, eb_color=VIOLET, ty=1.16, tw=CW,
           title_size=40, line=1.05):
    eyebrow(slide, M, 0.74, eb, color=eb_color)
    hline(slide, M, 1.04, 0.46, color=GOLD, weight=1.6)
    text(slide, M, ty, tw, 1.45, [{
        "line": line,
        "runs": title_runs,
    }])


# ---- card grid helper ------------------------------------------------------
def feature_card(slide, x, y, w, h, icon, title, desc,
                 icon_c=VIOLET, title_c=INK, desc_c=MUTED,
                 fill=PANEL, line=BORDER, accent_top=False, title_size=19):
    rrect(slide, x, y, w, h, fill=fill, line=line, line_w=1.0,
          shadow=True, shadow_kw={"opacity": 16, "blur": 0.10, "dist": 0.05})
    if accent_top:
        rect(slide, x + 0.0, y, 0.5, 0.055, fill=GOLD)
    pad = 0.32
    text(slide, x + pad, y + 0.32, w - 2 * pad, 0.7, [{
        "runs": [{"t": icon, "font": SYM, "s": 24, "c": icon_c}],
    }])
    text(slide, x + pad, y + 1.04, w - 2 * pad, 0.86, [{
        "line": 1.0,
        "runs": [{"t": title, "font": DISPLAY, "s": title_size, "c": title_c, "b": True}],
    }])
    text(slide, x + pad, y + 1.96, w - 2 * pad, h - 2.12, [{
        "line": 1.18,
        "runs": [{"t": desc, "font": BODY, "s": 11.2, "c": desc_c}],
    }])


# ============================================================================
# SLIDE 1 — COVER
# ============================================================================
def s_cover():
    s = slide_bg("bg_cream_cover.png")
    # top brand chip
    crown(s, M, 0.62, size=22, color=GOLD)
    text(s, M + 0.42, 0.66, 5, 0.4, [{
        "runs": [{"t": "ROYALCARDS", "font": BODY, "s": 12, "c": INK, "b": True,
                  "sp": 360}],
    }])
    text(s, SW - M - 5, 0.70, 5, 0.4, [{
        "align": PP_ALIGN.RIGHT,
        "runs": [{"t": "INVESTOR & PARTNER PITCH  ·  2026", "font": BODY,
                  "s": 10, "c": GOLD_DARK, "b": True, "sp": 200}],
    }])

    # left column — hero lockup
    lx = M
    text(s, lx, 1.74, 7.2, 1.5, [{
        "runs": [{"t": "RoyalCards", "font": SCRIPT, "s": 78, "c": VIOLET}],
    }])
    eyebrow(s, lx + 0.06, 3.18, "CURATED FOR COLLECTORS", color=GOLD_DARK)
    text(s, lx, 3.66, 6.6, 2.0, [
        {"line": 1.06, "runs": [
            {"t": "Königliche Schätze\n", "font": DISPLAY, "s": 52, "c": INK, "b": True}]},
        {"line": 1.06, "runs": [
            {"t": "für wahre ", "font": DISPLAY, "s": 52, "c": INK, "b": True},
            {"t": "Sammler", "font": DISPLAY, "s": 52, "c": VIOLET, "b": True, "i": True}]},
    ])
    text(s, lx, 5.66, 6.2, 1.0, [{
        "line": 1.25,
        "runs": [{"t": "Der kuratierte Premium-Marktplatz für seltene "
                  "TCG-Sammelkarten — authentifiziert, veredelt und mit "
                  "Stil präsentiert.", "font": BODY, "s": 12.5, "c": INK_SOFT}],
    }])

    # right column — fanned cards
    fan_cards(s)

    footer(s, 1)


def fan_cards(s):
    # (cx, cy, rot, fill, accent, rarity, name, sub, price, fill2)
    cards = [
        (9.05, 4.30, -13, VIOLET, GOLD_LIGHT, "◆ RARE", "Black Lotus",
         "Alpha · MTG", "25.000 €", VIOLET_DEEP),
        (11.05, 4.55, 9, "2E5E8C", "BFE0FF", "✦ SECRET RARE", "Blue-Eyes",
         "LOB · Yu-Gi-Oh!", "1.299 €", "1E3A5C"),
        (10.15, 2.55, -2, GOLD, "FFFFFF", "★ ULTRA RARE", "Charizard",
         "Base Set · 1999", "2.499 €", GOLD_DARK),
    ]
    cw, ch = 2.25, 3.25
    for (cx, cy, rot, fill, acc, rarity, name, sub, price, fill2) in cards:
        shp = slide_card(s, cx - cw / 2, cy - ch / 2, cw, ch, fill, fill2, rot)
        # text overlay (rotated with a separate box)
        tb = s.shapes.add_textbox(Inches(cx - cw / 2 + 0.16),
                                  Inches(cy - ch / 2 + 0.22),
                                  Inches(cw - 0.32), Inches(ch - 0.44))
        tb.rotation = rot
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = 0
        tf.margin_top = 0
        tf.margin_right = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        sym_part, _, txt_part = rarity.partition(" ")
        r0 = p.add_run()
        r0.text = sym_part + " "
        r0.font.name = SYM
        r0.font.size = Pt(8)
        r0.font.bold = True
        r0.font.color.rgb = C(acc)
        r = p.add_run()
        r.text = txt_part
        r.font.name = BODY
        r.font.size = Pt(8)
        r.font.bold = True
        r.font.color.rgb = C(acc)
        _set_spacing(r, 80)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(ch * 26)  # push name toward middle/bottom
        r2 = p2.add_run()
        r2.text = name
        r2.font.name = DISPLAY
        r2.font.size = Pt(20)
        r2.font.bold = True
        r2.font.color.rgb = C("FFFFFF")
        p3 = tf.add_paragraph()
        r3 = p3.add_run()
        r3.text = sub
        r3.font.name = BODY
        r3.font.size = Pt(8.5)
        r3.font.color.rgb = C(acc)
        p4 = tf.add_paragraph()
        p4.space_before = Pt(4)
        r4 = p4.add_run()
        r4.text = price
        r4.font.name = DISPLAY
        r4.font.size = Pt(17)
        r4.font.bold = True
        r4.font.color.rgb = C("FFFFFF")


def slide_card(s, x, y, w, h, fill, fill2, rot):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.08
    shp.rotation = rot
    try:
        shp.fill.gradient()
        stops = shp.fill.gradient_stops
        stops[0].position = 0.0
        stops[0].color.rgb = C(fill)
        stops[1].position = 1.0
        stops[1].color.rgb = C(fill2)
        shp.fill.gradient_angle = 65.0
    except Exception:
        shp.fill.solid()
        shp.fill.fore_color.rgb = C(fill)
    shp.line.color.rgb = C(GOLD_LIGHT)
    shp.line.width = Pt(1.25)
    no_shadow(shp)
    soft_shadow(shp, blur=0.16, dist=0.10, opacity=34, color="241433")
    return shp


# ============================================================================
# SLIDE 2 — VISION (violet statement)
# ============================================================================
def s_vision():
    s = slide_bg("bg_violet.png")
    crown(s, 0, 1.5, size=30, color=GOLD, align=PP_ALIGN.CENTER, w=SW)
    eyebrow(s, 0, 2.18, "DIE VISION", color=GOLD, w=SW, align=PP_ALIGN.CENTER)
    text(s, 1.2, 2.72, SW - 2.4, 2.3, [
        {"align": PP_ALIGN.CENTER, "line": 1.12, "runs": [
            {"t": "Wir machen das Sammeln von Trading Cards\n",
             "font": DISPLAY, "s": 37, "c": ON_V_TEXT, "b": True},
            {"t": "zu einem ", "font": DISPLAY, "s": 37, "c": ON_V_TEXT, "b": True},
            {"t": "königlichen Erlebnis", "font": DISPLAY, "s": 37,
             "c": GOLD, "b": True, "i": True},
            {"t": ".", "font": DISPLAY, "s": 37, "c": ON_V_TEXT, "b": True},
        ]},
    ], anchor=MSO_ANCHOR.TOP)
    hline(s, SW / 2 - 0.5, 5.18, 1.0, color=GOLD, weight=1.4)
    text(s, 2.4, 5.42, SW - 4.8, 0.8, [{
        "align": PP_ALIGN.CENTER, "line": 1.2,
        "runs": [{"t": "Kuratiert. Authentifiziert. Mit Stil präsentiert.  "
                  "Eleganz trifft Leidenschaft — seit 2014.",
                  "font": BODY, "s": 13, "c": ON_V_MUTED}],
    }])
    footer(s, 2, dark=True)


# ============================================================================
# SLIDE 3 — PROBLEM
# ============================================================================
def s_problem():
    s = slide_bg("bg_cream_content.png")
    header(s, "DAS PROBLEM", [
        {"t": "Ein boomender Markt — ", "font": DISPLAY, "s": 40, "c": INK, "b": True},
        {"t": "voller Unsicherheit", "font": DISPLAY, "s": 40, "c": VIOLET, "b": True, "i": True},
        {"t": ".", "font": DISPLAY, "s": 40, "c": INK, "b": True},
    ])
    cards = [
        ("⚠", "Fälschungen & Misstrauen",
         "Gefälschte und falsch deklarierte Karten kosten Sammler Geld — "
         "und dem Markt das Vertrauen."),
        ("◈", "Zersplitterter Markt",
         "Käufe verteilen sich auf eBay, Foren und Kleinanzeigen — ohne "
         "Garantie, ohne Kuratierung."),
        ("❖", "Kein Premium-Erlebnis",
         "Wertvolle Sammlerstücke verdienen mehr als einen anonymen, "
         "lieblosen Marktplatz."),
    ]
    gap = 0.36
    w = (CW - 2 * gap) / 3
    y = 2.55
    h = 3.55
    for i, (ic, t, d) in enumerate(cards):
        feature_card(s, M + i * (w + gap), y, w, h, ic, t, d,
                     icon_c=VIOLET, accent_top=True, title_size=20)
    footer(s, 3)


# ============================================================================
# SLIDE 4 — MARKT (violet emphasis)
# ============================================================================
def s_market():
    s = slide_bg("bg_violet_b.png")
    eyebrow(s, M, 0.74, "DER MARKT", color=GOLD)
    hline(s, M, 1.04, 0.46, color=GOLD, weight=1.6)
    text(s, M, 1.16, CW, 1.45, [{
        "runs": [
            {"t": "Sammelkarten sind eine ", "font": DISPLAY, "s": 38,
             "c": ON_V_TEXT, "b": True},
            {"t": "Anlageklasse", "font": DISPLAY, "s": 38, "c": GOLD,
             "b": True, "i": True},
            {"t": " geworden.", "font": DISPLAY, "s": 38, "c": ON_V_TEXT, "b": True},
        ],
    }])
    stats = [
        ("~ 9 Mrd. $", "Globaler TCG-Markt (2024)"),
        ("~ 8 % p.a.", "Erwartetes Marktwachstum"),
        ("~ 20 Mrd. $", "Marktprognose bis 2032"),
    ]
    gap = 0.36
    w = (CW - 2 * gap) / 3
    y = 2.5
    h = 2.0
    for i, (num, lab) in enumerate(stats):
        x = M + i * (w + gap)
        rrect(s, x, y, w, h, fill=ON_V_PANEL, line=ON_V_BORDER, line_w=1.0,
              shadow=True, shadow_kw={"opacity": 30, "color": "120A22",
                                      "blur": 0.12, "dist": 0.06})
        text(s, x, y + 0.42, w, 0.9, [{
            "align": PP_ALIGN.CENTER,
            "runs": [{"t": num, "font": DISPLAY, "s": 44, "c": GOLD, "b": True}],
        }])
        text(s, x, y + 1.34, w, 0.5, [{
            "align": PP_ALIGN.CENTER,
            "runs": [{"t": lab, "font": BODY, "s": 11, "c": ON_V_MUTED, "sp": 40}],
        }])
    text(s, M, 4.95, CW, 1.1, [{
        "align": PP_ALIGN.CENTER, "line": 1.32,
        "runs": [{"t": "Eine neue Sammlergeneration, Rekordpreise bei "
                  "Vintage-Karten, Nostalgie und Sammelkarten als "
                  "Wertanlage treiben die Nachfrage — quer durch alle "
                  "Spiele.", "font": BODY, "s": 13, "c": ON_V_TEXT}],
    }])
    text(s, M, 6.55, CW, 0.3, [{
        "align": PP_ALIGN.CENTER,
        "runs": [{"t": "Marktindikationen — vor finalem Pitch mit aktuellen "
                  "Quellen verifizieren.", "font": BODY, "s": 8.5,
                  "c": ON_V_MUTED, "i": True}],
    }])
    footer(s, 4, dark=True)


# ============================================================================
# SLIDE 5 — LÖSUNG
# ============================================================================
def s_solution():
    s = slide_bg("bg_cream_content.png")
    header(s, "DIE LÖSUNG", [
        {"t": "Vertrauen, Kuratierung, ", "font": DISPLAY, "s": 40, "c": INK, "b": True},
        {"t": "Eleganz", "font": DISPLAY, "s": 40, "c": VIOLET, "b": True, "i": True},
        {"t": ".", "font": DISPLAY, "s": 40, "c": INK, "b": True},
    ])
    pillars = [
        ("❖", "Kuratierte Auswahl",
         "Jede Karte handverlesen — keine Massenware, nur Stücke mit Charakter."),
        ("✔", "Geprüfte Authentizität",
         "Experten-Prüfung & Authentizitätsgarantie bei jedem einzelnen Kauf."),
        ("♛", "Grading inklusive",
         "PSA-, BGS- & CGC-Submission als Service — direkt aus Deutschland."),
        ("✦", "Königliches Erlebnis",
         "Fürstliche Verpackung, persönliche Beratung, Shop wie ein Schmuckstück."),
    ]
    gap = 0.30
    w = (CW - 3 * gap) / 4
    y = 2.7
    h = 3.4
    for i, (ic, t, d) in enumerate(pillars):
        feature_card(s, M + i * (w + gap), y, w, h, ic, t, d, accent_top=True,
                     title_size=18)
    footer(s, 5)


# ============================================================================
# SLIDE 6 — PLATTFORM (real screenshot)
# ============================================================================
def s_platform():
    s = slide_bg("bg_cream_content.png")
    eyebrow(s, M, 0.74, "DIE PLATTFORM", color=VIOLET)
    hline(s, M, 1.04, 0.46, color=GOLD, weight=1.6)
    text(s, M, 1.16, CW, 1.45, [{
        "runs": [
            {"t": "Ein Shop, gebaut wie ein ", "font": DISPLAY, "s": 38,
             "c": INK, "b": True},
            {"t": "Schmuckstück", "font": DISPLAY, "s": 38, "c": VIOLET,
             "b": True, "i": True},
            {"t": ".", "font": DISPLAY, "s": 38, "c": INK, "b": True},
        ],
    }])
    # screenshot framed (left)
    src = os.path.join(ASSETS, "site_hero.png")
    crop = os.path.join(ASSETS, "site_hero_crop.png")
    im = Image.open(src)
    iw, ih = im.size
    im.crop((0, 0, iw, int(ih * 0.905))).save(crop)
    cim = Image.open(crop)
    aspect = cim.size[0] / cim.size[1]
    pic_w = 6.9
    pic_h = pic_w / aspect
    px, py = M, 2.55
    # frame
    rrect(s, px - 0.10, py - 0.10, pic_w + 0.20, pic_h + 0.20, fill=WHITE,
          line=BORDER, line_w=1.0, radius=0.03,
          shadow=True, shadow_kw={"opacity": 24, "blur": 0.16, "dist": 0.10})
    pic = s.shapes.add_picture(crop, Inches(px), Inches(py),
                               width=Inches(pic_w), height=Inches(pic_h))
    pic.line.color.rgb = C(BORDER)
    pic.line.width = Pt(0.75)
    text(s, px, py + pic_h + 0.22, pic_w, 0.3, [{
        "align": PP_ALIGN.CENTER,
        "runs": [{"t": "royalcards.de", "font": BODY, "s": 9.5,
                  "c": MUTED, "sp": 150}],
    }])
    # features (right)
    fx = px + pic_w + 0.6
    fw = SW - M - fx
    feats = [
        ("Kuratierter Shop", "Stöbern nach Spiel & Produktart."),
        ("Wunschliste & Verlauf", "„Merken“ & „Zuletzt angesehen“."),
        ("Sicherer Checkout", "PayPal, Klarna, SEPA & Kreditkarte."),
        ("Exklusive Drops", "Seltene Karten zuerst — per Newsletter."),
        ("Mobil & blitzschnell", "Voll responsiv, Premium auf jedem Gerät."),
    ]
    fy = 2.62
    rowh = 0.84
    for i, (t, d) in enumerate(feats):
        yy = fy + i * rowh
        oval(s, fx, yy + 0.05, 0.13, 0.13, fill=GOLD)
        text(s, fx + 0.32, yy - 0.04, fw - 0.32, 0.76, [
            {"runs": [{"t": t, "font": DISPLAY, "s": 17, "c": INK, "b": True}]},
            {"space_before": 2, "line": 1.1, "runs": [
                {"t": d, "font": BODY, "s": 10.5, "c": MUTED}]},
        ])
    footer(s, 6)


# ============================================================================
# SLIDE 7 — SORTIMENT
# ============================================================================
def s_assortment():
    s = slide_bg("bg_cream_content.png")
    header(s, "DAS SORTIMENT", [
        {"t": "Sechs Welten. ", "font": DISPLAY, "s": 40, "c": INK, "b": True},
        {"t": "Eine Adresse", "font": DISPLAY, "s": 40, "c": VIOLET, "b": True, "i": True},
        {"t": ".", "font": DISPLAY, "s": 40, "c": INK, "b": True},
    ])
    games = [
        ("◉", "Pokémon", "4.200+ Karten"),
        ("✦", "Magic: The Gathering", "3.800+ Karten"),
        ("◆", "Yu-Gi-Oh!", "2.900+ Karten"),
        ("⚓", "One Piece TCG", "1.500+ Karten"),
        ("✧", "Disney Lorcana", "1.100+ Karten"),
        ("⬢", "Digimon", "900+ Karten"),
    ]
    gap = 0.30
    cols = 3
    w = (CW - (cols - 1) * gap) / cols
    h = 1.56
    y0 = 2.62
    vgap = 0.24
    for i, (ic, name, cnt) in enumerate(games):
        r, cidx = divmod(i, cols)
        x = M + cidx * (w + gap)
        y = y0 + r * (h + vgap)
        rrect(s, x, y, w, h, fill=PANEL, line=BORDER, line_w=1.0,
              shadow=True, shadow_kw={"opacity": 14, "blur": 0.09, "dist": 0.045})
        oval(s, x + 0.30, y + 0.34, 0.62, 0.62, fill="F1E9FB", line=BORDER_V, line_w=1.0)
        text(s, x + 0.30, y + 0.45, 0.62, 0.5, [{
            "align": PP_ALIGN.CENTER,
            "runs": [{"t": ic, "font": SYM, "s": 20, "c": VIOLET}],
        }])
        text(s, x + 1.10, y + 0.34, w - 1.3, 0.5, [{
            "runs": [{"t": name, "font": DISPLAY, "s": 18.5, "c": INK, "b": True}]}])
        text(s, x + 1.10, y + 0.86, w - 1.3, 0.4, [{
            "runs": [{"t": cnt, "font": BODY, "s": 10.5, "c": GOLD_DARK, "b": True,
                      "sp": 30}]}])
    text(s, M, 6.46, CW, 0.4, [{
        "align": PP_ALIGN.CENTER,
        "runs": [
            {"t": "Nach Produktart:   ", "font": BODY, "s": 11, "c": MUTED, "b": True, "sp": 40},
            {"t": "Singles   ·   Booster   ·   Displays   ·   Zubehör",
             "font": BODY, "s": 11, "c": INK, "sp": 40}],
    }])
    footer(s, 7)


# ============================================================================
# SLIDE 8 — WAS UNS AUSZEICHNET
# ============================================================================
def s_promise():
    s = slide_bg("bg_cream_content.png")
    header(s, "UNSER VERSPRECHEN", [
        {"t": "Vier Gründe, ", "font": DISPLAY, "s": 40, "c": INK, "b": True},
        {"t": "RoyalCards zu vertrauen", "font": DISPLAY, "s": 40, "c": VIOLET, "b": True, "i": True},
        {"t": ".", "font": DISPLAY, "s": 40, "c": INK, "b": True},
    ])
    props = [
        ("✔", "100 % Authentisch",
         "Jede Karte wird von Experten geprüft. Authentizitätsgarantie bei jedem Kauf."),
        ("✈", "Sicherer Versand",
         "Versichert, mit Sendungsverfolgung und professioneller Verpackung."),
        ("♛", "Grading-Service",
         "Professioneller PSA-, BGS- & CGC-Submission-Service aus Deutschland."),
        ("✉", "Experten-Support",
         "Persönliche Beratung durch echte TCG-Enthusiasten — auf Augenhöhe."),
    ]
    gap = 0.30
    w = (CW - 3 * gap) / 4
    y = 2.7
    h = 3.4
    for i, (ic, t, d) in enumerate(props):
        feature_card(s, M + i * (w + gap), y, w, h, ic, t, d, accent_top=True,
                     title_size=18)
    footer(s, 8)


# ============================================================================
# SLIDE 9 — VERTRAUEN / GRADING
# ============================================================================
def s_grading():
    s = slide_bg("bg_cream_content.png")
    header(s, "ZERTIFIZIERTE QUALITÄT", [
        {"t": "Bewertet von den ", "font": DISPLAY, "s": 40, "c": INK, "b": True},
        {"t": "Besten der Welt", "font": DISPLAY, "s": 40, "c": VIOLET, "b": True, "i": True},
        {"t": ".", "font": DISPLAY, "s": 40, "c": INK, "b": True},
    ])
    partners = [
        ("PSA", "Professional Sports Authenticator",
         "Weltweiter Marktführer im Card-Grading mit höchster Marktakzeptanz."),
        ("BGS", "Beckett Grading Services",
         "Bekannt für detaillierte Sub-Grades und die begehrten „Black Label“-Wertungen."),
        ("CGC", "Certified Guaranty Company",
         "Schnellste Bearbeitungszeiten bei ausgezeichneten Bewertungsdetails."),
    ]
    gap = 0.36
    w = (CW - 2 * gap) / 3
    y = 2.55
    h = 3.05
    for i, (abbr, name, d) in enumerate(partners):
        x = M + i * (w + gap)
        rrect(s, x, y, w, h, fill=PANEL, line=BORDER, line_w=1.0,
              shadow=True, shadow_kw={"opacity": 16, "blur": 0.10, "dist": 0.05})
        rect(s, x, y, 0.5, 0.055, fill=GOLD)
        # badge
        bw = 1.5
        rrect(s, x + (w - bw) / 2, y + 0.42, bw, 0.86, fill=INK, line=GOLD, line_w=1.25,
              radius=0.12, shadow=False)
        text(s, x + (w - bw) / 2, y + 0.58, bw, 0.6, [{
            "align": PP_ALIGN.CENTER,
            "runs": [{"t": abbr, "font": DISPLAY, "s": 30, "c": GOLD, "b": True, "sp": 80}],
        }])
        text(s, x + 0.3, y + 1.52, w - 0.6, 0.5, [{
            "align": PP_ALIGN.CENTER,
            "runs": [{"t": name, "font": DISPLAY, "s": 15.5, "c": INK, "b": True}]}])
        text(s, x + 0.34, y + 2.04, w - 0.68, 0.9, [{
            "align": PP_ALIGN.CENTER, "line": 1.18,
            "runs": [{"t": d, "font": BODY, "s": 11, "c": MUTED}]}])
    text(s, M, 5.92, CW, 0.4, [{
        "align": PP_ALIGN.CENTER,
        "runs": [{"t": "Professioneller Submission-Service — wir übernehmen "
                  "Einreichung, Versicherung und Rückversand.",
                  "font": BODY, "s": 11.5, "c": INK_SOFT, "i": True}],
    }])
    footer(s, 9)


# ============================================================================
# SLIDE 10 — TRAKTION
# ============================================================================
def s_traction():
    s = slide_bg("bg_cream_content.png")
    header(s, "TRAKTION", [
        {"t": "Über ", "font": DISPLAY, "s": 40, "c": INK, "b": True},
        {"t": "5.000 Sammler", "font": DISPLAY, "s": 40, "c": VIOLET, "b": True, "i": True},
        {"t": " vertrauen uns bereits.", "font": DISPLAY, "s": 40, "c": INK, "b": True},
    ])
    stats = [
        ("15.000+", "Karten im Sortiment", None),
        ("4.9", "Kundenbewertung", "★"),
        ("10+", "Jahre Erfahrung", None),
        ("5.000+", "Zufriedene Sammler", None),
    ]
    gap = 0.30
    w = (CW - 3 * gap) / 4
    y = 2.5
    h = 1.66
    for i, (num, lab, star) in enumerate(stats):
        x = M + i * (w + gap)
        rrect(s, x, y, w, h, fill=PANEL, line=BORDER, line_w=1.0,
              shadow=True, shadow_kw={"opacity": 14, "blur": 0.09, "dist": 0.045})
        num_runs = [{"t": num, "font": DISPLAY, "s": 34, "c": VIOLET, "b": True}]
        if star:
            num_runs.append({"t": " " + star, "font": SYM, "s": 24, "c": GOLD, "b": True})
        text(s, x, y + 0.30, w, 0.7, [{
            "align": PP_ALIGN.CENTER, "runs": num_runs}])
        text(s, x, y + 1.05, w, 0.4, [{
            "align": PP_ALIGN.CENTER,
            "runs": [{"t": lab, "font": BODY, "s": 10, "c": MUTED, "sp": 30}]}])
    # testimonial band
    ty = 4.55
    th = 1.85
    rrect(s, M, ty, CW, th, fill=CREAM_WARM, line=CREAM_DEEP, line_w=1.0,
          radius=0.04, shadow=True, shadow_kw={"opacity": 12, "blur": 0.10, "dist": 0.05})
    text(s, M + 0.5, ty + 0.22, 1.2, 1.0, [{
        "runs": [{"t": "“", "font": DISPLAY, "s": 80, "c": GOLD}]}], fit=False)
    text(s, M + 1.4, ty + 0.42, CW - 2.4, 1.1, [{
        "line": 1.22,
        "runs": [{"t": "Die Verpackung war fürstlich, die Karte makellos. "
                  "RoyalCards ist meine erste Adresse für seltene "
                  "Vintage-Karten geworden.", "font": DISPLAY, "s": 21,
                  "c": INK, "i": True}]}])
    text(s, M + 1.4, ty + th - 0.5, CW - 2.4, 0.4, [{
        "runs": [
            {"t": "VERIFIZIERTE KUNDENBEWERTUNG", "font": BODY, "s": 10, "c": INK, "b": True, "sp": 120},
            {"t": "   ·   RoyalCards", "font": BODY, "s": 10, "c": MUTED, "sp": 40}]}])
    footer(s, 10)


# ============================================================================
# SLIDE 11 — GESCHÄFTSMODELL
# ============================================================================
def s_business():
    s = slide_bg("bg_cream_content.png")
    header(s, "GESCHÄFTSMODELL", [
        {"t": "Mehrere ", "font": DISPLAY, "s": 40, "c": INK, "b": True},
        {"t": "Erlösströme", "font": DISPLAY, "s": 40, "c": VIOLET, "b": True, "i": True},
        {"t": ", eine Marke.", "font": DISPLAY, "s": 40, "c": INK, "b": True},
    ])
    streams = [
        ("◆", "Singles", "Margenstarker Verkauf einzelner Sammlerkarten."),
        ("■", "Sealed & Displays", "Versiegelte Booster-Boxen & Elite-Trainer-Boxen."),
        ("♛", "Grading-Service", "Service- & Vermittlungsgebühr auf PSA/BGS/CGC."),
        ("❖", "Zubehör", "Sleeves, Toploader, Alben & Playmats."),
    ]
    gap = 0.30
    w = (CW - 3 * gap) / 4
    y = 2.6
    h = 2.55
    for i, (ic, t, d) in enumerate(streams):
        x = M + i * (w + gap)
        rrect(s, x, y, w, h, fill=PANEL, line=BORDER, line_w=1.0,
              shadow=True, shadow_kw={"opacity": 15, "blur": 0.10, "dist": 0.05})
        text(s, x + 0.28, y + 0.28, w - 0.56, 0.6, [{
            "runs": [{"t": ic, "font": SYM, "s": 23, "c": GOLD_DARK}]}])
        text(s, x + 0.28, y + 0.92, w - 0.56, 0.5, [{
            "runs": [{"t": t, "font": DISPLAY, "s": 18, "c": INK, "b": True}]}])
        text(s, x + 0.28, y + 1.42, w - 0.56, 1.0, [{
            "line": 1.16,
            "runs": [{"t": d, "font": BODY, "s": 10.5, "c": MUTED}]}])
    # future band
    by = 5.45
    bh = 0.92
    rrect(s, M, by, CW, bh, fill=INK, line=None, radius=0.10,
          shadow=True, shadow_kw={"opacity": 22, "blur": 0.12, "dist": 0.06})
    text(s, M + 0.5, by, CW - 1.0, bh, [{
        "runs": [
            {"t": "♛", "font": SYM, "s": 13, "c": GOLD, "b": True},
            {"t": "  ZUKUNFT   ", "font": BODY, "s": 11, "c": GOLD, "b": True, "sp": 160},
            {"t": "Consignment & Auktionen — wir öffnen RoyalCards als "
             "kuratierten Marktplatz für Verkäufer.",
             "font": BODY, "s": 12.5, "c": ON_V_TEXT}],
    }], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 11)


# ============================================================================
# SLIDE 12 — GO-TO-MARKET / ROADMAP
# ============================================================================
def s_roadmap():
    s = slide_bg("bg_cream_content.png")
    header(s, "FAHRPLAN", [
        {"t": "Der Weg ", "font": DISPLAY, "s": 40, "c": INK, "b": True},
        {"t": "an den Markt", "font": DISPLAY, "s": 40, "c": VIOLET, "b": True, "i": True},
        {"t": ".", "font": DISPLAY, "s": 40, "c": INK, "b": True},
    ])
    phases = [
        ("01", "Launch", "Sortimentstiefe aufbauen, Shop-Erlebnis perfektionieren."),
        ("02", "Community", "Newsletter-Drops, Social Media & lokale Sammler-Events."),
        ("03", "Grading skalieren", "Partnerschaften vertiefen, Durchlaufzeiten senken."),
        ("04", "Marktplatz", "Consignment & Auktionen öffnen — Angebot vervielfachen."),
    ]
    gap = 0.30
    w = (CW - 3 * gap) / 4
    y = 3.0
    h = 2.9
    # connecting line
    hline(s, M + w / 2, y + 0.05, CW - w, color=CREAM_DEEP, weight=2.0)
    for i, (num, t, d) in enumerate(phases):
        x = M + i * (w + gap)
        # node
        nd = 0.62
        oval(s, x + w / 2 - nd / 2, y - nd / 2, nd, nd, fill=VIOLET, line=GOLD, line_w=1.5)
        text(s, x + w / 2 - nd / 2, y - nd / 2 + 0.12, nd, 0.4, [{
            "align": PP_ALIGN.CENTER,
            "runs": [{"t": num, "font": DISPLAY, "s": 19, "c": "FFFFFF", "b": True}]}])
        # card
        cy = y + 0.62
        rrect(s, x, cy, w, h - 0.62, fill=PANEL, line=BORDER, line_w=1.0,
              shadow=True, shadow_kw={"opacity": 14, "blur": 0.09, "dist": 0.045})
        text(s, x + 0.26, cy + 0.30, w - 0.52, 0.5, [{
            "align": PP_ALIGN.CENTER,
            "runs": [{"t": t, "font": DISPLAY, "s": 18, "c": INK, "b": True}]}])
        text(s, x + 0.26, cy + 0.86, w - 0.52, 1.2, [{
            "align": PP_ALIGN.CENTER, "line": 1.18,
            "runs": [{"t": d, "font": BODY, "s": 10.5, "c": MUTED}]}])
    footer(s, 12)


# ============================================================================
# SLIDE 13 — CLOSING / ASK
# ============================================================================
def s_closing():
    s = slide_bg("bg_violet_c.png")
    crown(s, 0, 1.35, size=32, color=GOLD, align=PP_ALIGN.CENTER, w=SW)
    eyebrow(s, 0, 2.05, "LASSEN SIE UNS SPRECHEN", color=GOLD, w=SW, align=PP_ALIGN.CENTER)
    text(s, 1.3, 2.58, SW - 2.6, 1.7, [
        {"align": PP_ALIGN.CENTER, "line": 1.12, "runs": [
            {"t": "Werden Sie Teil der\n", "font": DISPLAY, "s": 40, "c": ON_V_TEXT, "b": True},
            {"t": "königlichen Kollektion", "font": DISPLAY, "s": 40, "c": GOLD, "b": True, "i": True},
            {"t": ".", "font": DISPLAY, "s": 40, "c": ON_V_TEXT, "b": True}]},
    ])
    text(s, 2.6, 4.35, SW - 5.2, 0.9, [{
        "align": PP_ALIGN.CENTER, "line": 1.28,
        "runs": [{"t": "Wir suchen Partner und Investoren, um RoyalCards zur "
                  "ersten Adresse für Premium-TCG im DACH-Raum zu machen.",
                  "font": BODY, "s": 13.5, "c": ON_V_MUTED}],
    }])
    hline(s, SW / 2 - 0.6, 5.5, 1.2, color=GOLD, weight=1.2)
    # contact line
    text(s, 0, 5.74, SW, 0.4, [{
        "align": PP_ALIGN.CENTER,
        "runs": [
            {"t": "RoyalCards GmbH", "font": DISPLAY, "s": 17, "c": ON_V_TEXT, "b": True},
            {"t": "      Königsallee 42, 40212 Düsseldorf", "font": BODY, "s": 11.5, "c": ON_V_MUTED}],
    }])
    text(s, 0, 6.18, SW, 0.4, [{
        "align": PP_ALIGN.CENTER,
        "runs": [
            {"t": "royalcards.de", "font": BODY, "s": 11.5, "c": GOLD, "b": True, "sp": 60},
            {"t": "      ·      ", "font": BODY, "s": 11.5, "c": ON_V_MUTED},
            {"t": "kontakt@royalcards.de", "font": BODY, "s": 11.5, "c": GOLD, "b": True, "sp": 60}],
    }])
    footer(s, 13, dark=True)


# ============================================================================
# MOTION — gentle fade transition + auto, staggered fade-in builds
# ============================================================================
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as _AST  # noqa: E402


def _area_in(sh):
    try:
        return (sh.width / 914400.0) * (sh.height / 914400.0)
    except Exception:
        return 0.0


def _is_container(sh):
    try:
        a = sh.auto_shape_type
    except Exception:
        return False
    if a not in (_AST.ROUNDED_RECTANGLE, _AST.RECTANGLE, _AST.OVAL):
        return False
    return _area_in(sh) >= 1.5


def _groups(slide):
    """Cluster shapes into build groups (a panel + its contents = one beat)."""
    sw = prs.slide_width
    groups = []
    cur = None
    cont = None
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and sh.width >= sw * 0.97:
            continue  # full-bleed background — always visible
        if sh.top is not None and sh.top >= Inches(6.85):
            continue  # footer — always visible
        sid = sh.shape_id
        if _is_container(sh):
            cur = [sid]
            groups.append(cur)
            cont = (sh.left, sh.top, sh.width, sh.height)
        else:
            if cont is not None:
                cx = sh.left + sh.width / 2
                cy = sh.top + sh.height / 2
                bx, by, bw, bh = cont
                if (bx - 9144 <= cx <= bx + bw + 9144
                        and by - 9144 <= cy <= by + bh + 9144):
                    cur.append(sid)
                    continue
            cur = [sid]
            groups.append(cur)
            cont = None
    return groups


def _eff_par(cid, sid, node, delay, dur, grp):
    c1, c2, c3 = cid, cid + 1, cid + 2
    return ('<p:par><p:cTn id="%d" presetID="10" presetClass="entr" '
            'presetSubtype="0" fill="hold" grpId="%d" nodeType="%s">'
            '<p:stCondLst><p:cond delay="%d"/></p:stCondLst><p:childTnLst>'
            '<p:set><p:cBhvr><p:cTn id="%d" dur="1" fill="hold">'
            '<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
            '<p:tgtEl><p:spTgt spid="%d"/></p:tgtEl>'
            '<p:attrNameLst><p:attrName>style.visibility</p:attrName>'
            '</p:attrNameLst></p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
            '<p:animEffect transition="in" filter="fade"><p:cBhvr>'
            '<p:cTn id="%d" dur="%d"/><p:tgtEl><p:spTgt spid="%d"/></p:tgtEl>'
            '</p:cBhvr></p:animEffect></p:childTnLst></p:cTn></p:par>'
            % (c1, grp, node, delay, c2, sid, c3, dur, sid))


def _timing_xml(groups, dur=480, cascade=210, initial=200):
    cid = 100
    pars = []
    first = True
    for gi, g in enumerate(groups):
        for si, sid in enumerate(g):
            if first:
                node, delay = "afterEffect", initial
                first = False
            elif si == 0:
                node, delay = "withEffect", cascade
            else:
                node, delay = "withEffect", 0
            pars.append(_eff_par(cid, sid, node, delay, dur, gi))
            cid += 3
    body = "".join(pars)
    return (
        '<p:timing xmlns:p="%s" xmlns:a="%s"><p:tnLst><p:par>'
        '<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
        '<p:childTnLst><p:seq concurrent="1" nextAc="seek">'
        '<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>%s'
        '</p:childTnLst></p:cTn>'
        '<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/>'
        '</p:tgtEl></p:cond></p:prevCondLst>'
        '<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/>'
        '</p:tgtEl></p:cond></p:nextCondLst>'
        '</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>'
        % (P_NS, A_NS, body))


def add_motion(slide):
    el = slide._element
    el.append(etree.fromstring(
        '<p:transition xmlns:p="%s" spd="med"><p:fade/></p:transition>' % P_NS))
    groups = _groups(slide)
    if groups:
        el.append(etree.fromstring(_timing_xml(groups)))


# ----------------------------------------------------------------------------
for fn in (s_cover, s_vision, s_problem, s_market, s_solution, s_platform,
           s_assortment, s_promise, s_grading, s_traction, s_business,
           s_roadmap, s_closing):
    fn()

for _slide in prs.slides:
    add_motion(_slide)

prs.save(OUT)
print("saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
