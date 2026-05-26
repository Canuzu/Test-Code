# -*- coding: utf-8 -*-
"""Approximate PPTX -> PNG renderer (QA only).

Reads the real .pptx via python-pptx and rasterises each slide with PIL using
the installed brand fonts. Not pixel-perfect vs PowerPoint, but faithful enough
to catch layout problems (overflow, overlaps, off-canvas, the rotated cards).
"""
import os
import math
import numpy as np
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as AST
from pptx.oxml.ns import qn

HERE = os.path.dirname(__file__)
PPTX = os.path.join(os.path.dirname(HERE), "RoyalCards_Pitch_Deck.pptx")
OUTDIR = os.path.join(HERE, "render")
os.makedirs(OUTDIR, exist_ok=True)

SCALE = 150.0  # px per inch
EMU_IN = 914400.0
DEBUG = os.environ.get("DEBUG_BOXES") == "1"
WARN = []
_CUR = {"slide": 0}
FB = "/usr/share/fonts/truetype/royalbrand"
DEJAVU = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
if not os.path.exists(DEJAVU):
    import subprocess
    DEJAVU = subprocess.run(["fc-match", "-f", "%{file}", "DejaVu Sans"],
                            capture_output=True, text=True).stdout.strip() or None

FONT_FILES = {
    "Cormorant Garamond": (f"{FB}/CormorantGaramond[wght].ttf",
                           f"{FB}/CormorantGaramond-Italic[wght].ttf"),
    "Great Vibes": (f"{FB}/GreatVibes-Regular.ttf", f"{FB}/GreatVibes-Regular.ttf"),
    "Inter": (f"{FB}/Inter[opsz,wght].ttf", f"{FB}/Inter-Italic[opsz,wght].ttf"),
}
_fc = {}


def px(emu):
    return emu / EMU_IN * SCALE


def get_font(name, size_pt, bold, italic):
    size = max(4, int(round(size_pt * SCALE / 72.0)))
    key = (name, size, bool(bold), bool(italic))
    if key in _fc:
        return _fc[key]
    path = None
    sym = False
    if name in FONT_FILES:
        path = FONT_FILES[name][1 if italic else 0]
    else:  # symbol / unknown -> DejaVu
        path = DEJAVU
        sym = True
    try:
        f = ImageFont.truetype(path, size)
    except Exception:
        f = ImageFont.truetype(DEJAVU, size)
        sym = True
    if not sym:
        try:
            target = "Bold" if bold else "Regular"
            names = [n.decode() if isinstance(n, bytes) else n
                     for n in f.get_variation_names()]
            if target in names:
                f.set_variation_by_name(target)
            elif bold and "Bold" in names:
                f.set_variation_by_name("Bold")
        except Exception:
            pass
    _fc[key] = f
    return f


def run_props(run):
    f = run.font
    name = f.name or "Inter"
    size = f.size.pt if f.size is not None else 14
    bold = bool(f.bold)
    italic = bool(f.italic)
    try:
        col = f.color.rgb
        color = (col[0], col[1], col[2])
    except Exception:
        color = (40, 26, 77)
    spc = 0
    rPr = run._r.find(qn("a:rPr"))
    if rPr is not None and rPr.get("spc"):
        spc = int(rPr.get("spc")) / 100.0  # pt
    return dict(text=run.text, name=name, size=size, bold=bold,
                italic=italic, color=color, spc=spc)


def measure(draw, txt, font, spc_px):
    if spc_px <= 0:
        return draw.textlength(txt, font=font)
    return sum(draw.textlength(ch, font=font) for ch in txt) + spc_px * max(0, len(txt) - 1)


def draw_run_text(draw, x, y, txt, font, color, spc_px):
    if spc_px <= 0:
        draw.text((x, y), txt, font=font, fill=color)
        return x + draw.textlength(txt, font=font)
    for ch in txt:
        draw.text((x, y), ch, font=font, fill=color)
        x += draw.textlength(ch, font=font) + spc_px
    return x


def layout_paragraphs(draw, box, tf, vanchor):
    """box=(x,y,w,h) px. Returns list of lines; each line list of segments."""
    bx, by, bw, bh = box
    lines = []  # (segments, height, ascent) segments:(txt,font,color,spc_px,width)
    for para in tf.paragraphs:
        runs = [run_props(r) for r in para.runs if r.text != ""]
        align = para.alignment or PP_ALIGN.LEFT
        ls = para.line_spacing
        sb = para.space_before.pt * SCALE / 72.0 if para.space_before else 0
        sa = para.space_after.pt * SCALE / 72.0 if para.space_after else 0
        if not runs:
            lines.append(dict(segs=[], h=10, align=align, sb=sb, sa=sa))
            continue
        # tokenize into words keeping run styling (handle forced \n breaks)
        tokens = []  # (word, props) ; word == "\n" => forced break
        for rp in runs:
            segs = rp["text"].split("\n")
            for si, seg in enumerate(segs):
                if si > 0:
                    tokens.append(("\n", rp))
                parts = seg.split(" ")
                for j, w in enumerate(parts):
                    if j > 0:
                        tokens.append((" ", rp))
                    if w != "":
                        tokens.append((w, rp))
        cur = []
        cur_w = 0
        max_h = 0
        first = True

        def flush():
            nonlocal cur, cur_w, max_h, first
            lines.append(dict(segs=cur, h=max_h, align=align,
                              sb=(sb if first else 0), sa=0))
            cur = []
            cur_w = 0
            max_h = 0
            first = False
        for word, rp in tokens:
            font = get_font(rp["name"], rp["size"], rp["bold"], rp["italic"])
            spc_px = rp["spc"] * SCALE / 72.0
            if word == "\n":
                if max_h == 0:
                    asc, desc = font.getmetrics()
                    max_h = asc + desc
                flush()
                continue
            ww = measure(draw, word, font, spc_px)
            asc, desc = font.getmetrics()
            hh = (asc + desc)
            if word == " ":
                if not cur:
                    continue
                cur.append([" ", font, rp["color"], spc_px, ww])
                cur_w += ww
                max_h = max(max_h, hh)
                continue
            if cur and cur_w + ww > bw:
                # strip trailing space
                while cur and cur[-1][0] == " ":
                    cur_w -= cur[-1][4]
                    cur.pop()
                flush()
            cur.append([word, font, rp["color"], spc_px, ww])
            cur_w += ww
            max_h = max(max_h, hh)
        while cur and cur[-1][0] == " ":
            cur.pop()
        flush()
        # apply space_after to last line of paragraph
        if lines:
            lines[-1]["sa"] = sa
            lines[-1]["ls"] = ls
        for ln in lines:
            ln.setdefault("ls", ls)
    # compute total height
    total = 0
    for ln in lines:
        lh = ln["h"] * (ln["ls"] if ln["ls"] else 1.0)
        total += ln["sb"] + lh + ln["sa"]
    if vanchor == MSO_ANCHOR.MIDDLE:
        cy = by + (bh - total) / 2
    elif vanchor == MSO_ANCHOR.BOTTOM:
        cy = by + (bh - total)
    else:
        cy = by
    # draw
    for ln in lines:
        lh = ln["h"] * (ln["ls"] if ln["ls"] else 1.0)
        cy += ln["sb"]
        line_w = sum(seg[4] for seg in ln["segs"])
        if ln["align"] == PP_ALIGN.CENTER:
            cx = bx + (bw - line_w) / 2
        elif ln["align"] == PP_ALIGN.RIGHT:
            cx = bx + (bw - line_w)
        else:
            cx = bx
        # baseline padding inside line height
        pad = (lh - ln["h"]) / 2
        for (txt, font, color, spc_px, w) in ln["segs"]:
            draw_run_text(draw, cx, cy + pad, txt, font, color, spc_px)
            cx += w
        cy += lh + ln["sa"]
    return total


def shape_fill(shape):
    try:
        ft = shape.fill.type
    except Exception:
        return None
    if ft == 1:  # solid
        try:
            c = shape.fill.fore_color.rgb
            return ("solid", (c[0], c[1], c[2]))
        except Exception:
            return None
    if ft == 3:  # gradient
        try:
            stops = []
            for gs in shape.fill.gradient_stops:
                c = gs.color.rgb
                stops.append((gs.position, (c[0], c[1], c[2])))
            ang = 90
            try:
                ang = shape.fill.gradient_angle
            except Exception:
                pass
            return ("grad", stops, ang)
        except Exception:
            return None
    return None


def shape_line(shape):
    try:
        lt = shape.line.fill.type
        if lt is None or lt == 5:  # none/background
            return None
        c = shape.line.color.rgb
        w = shape.line.width.pt if shape.line.width else 1.0
        return ((c[0], c[1], c[2]), max(1, int(round(w * SCALE / 72.0))))
    except Exception:
        return None


def rounded_mask(w, h, rad):
    m = Image.new("L", (w, h), 0)
    d = ImageDraw.Draw(m)
    d.rounded_rectangle([0, 0, w - 1, h - 1], radius=rad, fill=255)
    return m


def grad_image(w, h, stops, angle):
    if len(stops) < 2:
        col = stops[0][1] if stops else (128, 128, 128)
        return Image.new("RGB", (w, h), col)
    c0 = np.array(stops[0][1], float)
    c1 = np.array(stops[-1][1], float)
    a = math.radians(angle)
    yy, xx = np.mgrid[0:h, 0:w].astype(float)
    nx, ny = math.cos(a), math.sin(a)
    t = (xx * nx + yy * ny)
    t = (t - t.min()) / (t.max() - t.min() + 1e-6)
    img = c0[None, None, :] * (1 - t[:, :, None]) + c1[None, None, :] * t[:, :, None]
    return Image.fromarray(np.clip(img, 0, 255).astype(np.uint8), "RGB")


def auto_kind(shape):
    try:
        a = shape.auto_shape_type
    except Exception:
        return None
    if a == AST.ROUNDED_RECTANGLE:
        return "rrect"
    if a == AST.OVAL:
        return "oval"
    if a == AST.RECTANGLE:
        return "rect"
    return "rect"


def render_shape(canvas, shape):
    # picture?
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            from io import BytesIO
            im = Image.open(BytesIO(shape.image.blob)).convert("RGB")
            w, h = int(px(shape.width)), int(px(shape.height))
            im = im.resize((max(1, w), max(1, h)))
            canvas.paste(im, (int(px(shape.left)), int(px(shape.top))))
            ln = shape_line(shape)
            if ln:
                d = ImageDraw.Draw(canvas)
                d.rectangle([px(shape.left), px(shape.top),
                             px(shape.left) + w, px(shape.top) + h],
                            outline=ln[0], width=ln[1])
        except Exception as e:
            print("pic err", e)
        return

    x, y = px(shape.left), px(shape.top)
    w, h = px(shape.width), px(shape.height)
    rot = float(shape.rotation or 0)
    has_text = shape.has_text_frame and shape.text_frame.text.strip() != ""
    fill = shape_fill(shape)
    line = shape_line(shape)
    kind = auto_kind(shape)

    # connector / zero-dim => line
    if (shape.width == 0 or shape.height == 0) and line:
        d = ImageDraw.Draw(canvas)
        d.line([x, y, x + w, y + h], fill=line[0], width=line[1])
        return

    iw, ih = max(1, int(round(w))), max(1, int(round(h)))
    tile = Image.new("RGBA", (iw, ih), (0, 0, 0, 0))
    td = ImageDraw.Draw(tile)
    rad = 0
    if kind == "rrect":
        try:
            rad = int(shape.adjustments[0] * min(iw, ih))
        except Exception:
            rad = int(0.05 * min(iw, ih))

    # fill
    if fill:
        if fill[0] == "solid":
            base = Image.new("RGBA", (iw, ih), fill[1] + (255,))
        else:
            base = grad_image(iw, ih, fill[1], fill[2]).convert("RGBA")
        if kind == "oval":
            m = Image.new("L", (iw, ih), 0)
            ImageDraw.Draw(m).ellipse([0, 0, iw - 1, ih - 1], fill=255)
        elif kind == "rrect":
            m = rounded_mask(iw, ih, rad)
        else:
            m = Image.new("L", (iw, ih), 255)
        tile.paste(base, (0, 0), m)
    # outline
    if line:
        if kind == "oval":
            td.ellipse([line[1] / 2, line[1] / 2, iw - 1 - line[1] / 2, ih - 1 - line[1] / 2],
                       outline=line[0], width=line[1])
        elif kind == "rrect":
            td.rounded_rectangle([line[1] / 2, line[1] / 2, iw - 1 - line[1] / 2, ih - 1 - line[1] / 2],
                                 radius=rad, outline=line[0], width=line[1])
        else:
            td.rectangle([line[1] / 2, line[1] / 2, iw - 1 - line[1] / 2, ih - 1 - line[1] / 2],
                         outline=line[0], width=line[1])
    # text
    if has_text:
        tf = shape.text_frame
        ml = px(tf.margin_left or 0)
        mr = px(tf.margin_right or 0)
        mt = px(tf.margin_top or 0)
        mb = px(tf.margin_bottom or 0)
        box = (ml, mt, iw - ml - mr, ih - mt - mb)
        total_h = layout_paragraphs(td, box, tf, tf.vertical_anchor)
        if total_h > box[3] + 6:  # text taller than box -> overflow risk
            snippet = tf.text.replace("\n", " / ")[:46]
            WARN.append("slide %02d OVERFLOW +%.2f\" : %s"
                        % (_CUR["slide"], (total_h - box[3]) / SCALE, snippet))
            if DEBUG:
                td.rectangle([0, 0, iw - 1, ih - 1], outline=(220, 40, 40), width=2)
        elif DEBUG:
            td.rectangle([0, 0, iw - 1, ih - 1], outline=(40, 140, 220), width=1)

    if abs(rot) > 0.01:
        tile = tile.rotate(-rot, expand=True, resample=Image.BICUBIC)
        cx, cy = x + w / 2, y + h / 2
        canvas.paste(tile, (int(round(cx - tile.size[0] / 2)),
                            int(round(cy - tile.size[1] / 2))), tile)
    else:
        canvas.paste(tile, (int(round(x)), int(round(y))), tile)


def render():
    prs = Presentation(PPTX)
    W = int(px(prs.slide_width))
    H = int(px(prs.slide_height))
    paths = []
    for i, slide in enumerate(prs.slides, 1):
        _CUR["slide"] = i
        canvas = Image.new("RGB", (W, H), (250, 246, 239))
        for shape in slide.shapes:
            render_shape(canvas, shape)
        p = os.path.join(OUTDIR, f"slide_{i:02d}.png")
        canvas.save(p)
        paths.append(p)
        print("rendered", p)
    # contact sheet
    cols, rows = 3, 5
    tw, th = W // 3, H // 3
    sheet = Image.new("RGB", (cols * tw + 40, rows * th + 60), (24, 16, 40))
    for idx, p in enumerate(paths):
        im = Image.open(p).resize((tw - 12, th - 12))
        r, c = divmod(idx, cols)
        sheet.paste(im, (20 + c * tw, 20 + r * th))
    sheet.save(os.path.join(OUTDIR, "contact_sheet.png"))
    print("contact sheet saved")
    print("\n=== OVERFLOW REPORT (%d) ===" % len(WARN))
    for w in WARN:
        print(" ", w)


if __name__ == "__main__":
    render()
